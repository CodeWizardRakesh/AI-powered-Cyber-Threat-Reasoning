from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import matplotlib.pyplot as plt
import os

# Initialize presentation
prs = Presentation()

# Slide content
slides = [
    ("Introduction", [
        "Cybersecurity is critical due to evolving threats like brute-force attacks and malware.",
        '"AI-Powered Cyber Threat Reasoning" uses AI to analyze logs and generate reports.',
        "Built with Flask and Gemini 1.5 Flash for real-time insights and downloads.",
        "Aims to empower users with actionable security recommendations."
    ]),
    ("Methodology", [
        "Defined goal: Analyze logs, generate reports, and enable downloads.",
        "Selected tools: Python, Flask, Gemini, Pandas, python-docx.",
        "Designed architecture: File upload, AI analysis, streaming, and export.",
        "Prepared resources: API key, sample logs, environment setup.",
        "Planned timeline: 4 weeks for development and testing."
    ]),
    ("Result", [
        "Technical: Python, web dev, API integration, doc generation.",
        "Analytical: Data processing, pattern recognition, debugging.",
        "AI/NLP: Prompt crafting, response parsing with regex.",
        "Design: UI/UX, real-time streaming.",
        "Problem-Solving: Adapted to formats, added downloads."
    ]),
    ("Output", [
        "Top Threat Sources: IP frequency list.",
        "High-Risk Events: Filtered log details.",
        "Report: Summary, patterns, actions streamed in UI.",
        "Download: Word doc with structured report."
    ]),
    ("Conclusion", [
        "Successfully built an AI-driven cybersecurity tool.",
        "Enhanced skills in coding, AI, and UX design.",
        "Offers practical value for threat analysis.",
        "Future scope: Multi-user support, advanced AI models."
    ])
]

# Function to create a bar chart image with Matplotlib
def create_threat_source_chart():
    # Sample data from your output: {'192.168.1.10': 1, '192.168.1.15': 1, '203.0.113.5': 1}
    ips = ['192.168.1.10', '192.168.1.15', '203.0.113.5']
    counts = [1, 1, 1]  # Replace with your actual data if different
    plt.bar(ips, counts, color='green')
    plt.title("Top Threat Sources")
    plt.xlabel("Source IP")
    plt.ylabel("Event Count")
    plt.xticks(rotation=45)
    plt.tight_layout()
    chart_file = "threat_sources_chart.png"
    plt.savefig(chart_file)
    plt.close()
    return chart_file

# Generate all slides
for i, (title, bullets) in enumerate(slides):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet

    # Add visualizations to the "Output" slide (Slide 4, index 3)
    if i == 3:  # "Output" slide
        # Add Bar Chart
        chart_file = create_threat_source_chart()
        slide.shapes.add_picture(chart_file, Inches(4), Inches(1.5), width=Inches(3), height=Inches(2))

        # Add Placeholder for Screenshot (replace 'screenshot.png' with your actual file)
        screenshot_path = "screenshot.png"  # Provide your app screenshot here
        if os.path.exists(screenshot_path):
            slide.shapes.add_picture(screenshot_path, Inches(4), Inches(4), width=Inches(3), height=Inches(2))
        else:
            print(f"Warning: '{screenshot_path}' not found. Add your screenshot to include it.")

# Save the presentation
prs.save("Cyber_Threat_Reasoning_With_Visuals1.pptx")
print("Presentation saved as 'Cyber_Threat_Reasoning_With_Visuals.pptx'")